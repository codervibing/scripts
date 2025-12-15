def overlap_ratio_with_canvas(e) -> float:
    """
    Returns fraction of the element's area that lies inside the slide canvas.
    Canvas is normalized [0,1] x [0,1].
    """
    b = e["bbox"]["norm"]
    x1, y1 = b["x"], b["y"]
    x2, y2 = b["x"] + b["w"], b["y"] + b["h"]

    # intersection with [0,1]x[0,1]
    ix1, iy1 = max(0.0, x1), max(0.0, y1)
    ix2, iy2 = min(1.0, x2), min(1.0, y2)

    iw, ih = max(0.0, ix2 - ix1), max(0.0, iy2 - iy1)
    inter = iw * ih
    area = max(1e-12, (x2 - x1) * (y2 - y1))
    return inter / area


def is_on_canvas(e, min_overlap: float = 0.90) -> bool:
    """
    Keep only elements where >= min_overlap of their area is inside the slide canvas.
    min_overlap=0.90 means drop elements mostly outside.
    """
    return overlap_ratio_with_canvas(e) >= min_overlap


#!/usr/bin/env python3
"""
pptx_structured_extract.py

Extracts a PPTX into a structured JSON representation WITHOUT any vision/OCR.
Works best when slides are made of native PPT objects (text boxes, tables, charts, shapes).
If content is pasted as images, you'll only get image bounding boxes + filenames.

Dependencies:
  pip install python-pptx

Usage:
  python pptx_structured_extract.py input.pptx -o out.json
"""

import argparse
import json
import math
import os
import re
import zipfile
from collections import defaultdict
from dataclasses import dataclass, asdict
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

import xml.etree.ElementTree as ET


EMU_PER_INCH = 914400


# ----------------------------
# Helpers: geometry + safe access
# ----------------------------
def emu_box(shape, slide_w: int, slide_h: int) -> Dict[str, Any]:
    """Return bbox in both EMU and normalized [0..1] coords."""
    x = int(getattr(shape, "left", 0) or 0)
    y = int(getattr(shape, "top", 0) or 0)
    w = int(getattr(shape, "width", 0) or 0)
    h = int(getattr(shape, "height", 0) or 0)
    nx = x / slide_w if slide_w else 0.0
    ny = y / slide_h if slide_h else 0.0
    nw = w / slide_w if slide_w else 0.0
    nh = h / slide_h if slide_h else 0.0
    return {
        "emu": {"x": x, "y": y, "w": w, "h": h},
        "norm": {"x": nx, "y": ny, "w": nw, "h": nh},
        "center_norm": {"x": nx + nw / 2.0, "y": ny + nh / 2.0},
    }


def safe_text(shape) -> str:
    try:
        if shape.has_text_frame and shape.text_frame:
            return shape.text_frame.text or ""
    except Exception:
        pass
    return ""


def max_font_pt(shape) -> Optional[float]:
    """Approximate the max run font size within a shape (points)."""
    try:
        if not shape.has_text_frame or not shape.text_frame:
            return None
        mx = None
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                if run.font.size is not None:
                    pt = float(run.font.size.pt)
                    mx = pt if mx is None else max(mx, pt)
        return mx
    except Exception:
        return None


def shape_type_name(shape) -> str:
    st = getattr(shape, "shape_type", None)
    if st == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if st == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if st == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if st == MSO_SHAPE_TYPE.PLACEHOLDER:
        # could be title/body etc
        return "placeholder"
    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        # could be textbox too
        if getattr(shape, "has_text_frame", False):
            return "textbox"
        return "shape"
    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        return "textbox"
    if st == MSO_SHAPE_TYPE.LINE:
        return "line"
    # fallback
    if getattr(shape, "has_text_frame", False):
        return "textbox"
    return "unknown"


# ----------------------------
# Chart XML parsing (OOXML) without vision
# ----------------------------
NS = {
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}

def _strip_pack_uri(pack_uri: str) -> str:
    """
    python-pptx chart partname looks like '/ppt/charts/chart1.xml'
    In zipfile, names don't start with leading slash.
    """
    return pack_uri.lstrip("/")


def parse_chart_xml_from_zip(zf: zipfile.ZipFile, chart_partname: str) -> Dict[str, Any]:
    """
    Extract series/categories/values from chart XML if present.

    Notes:
      - Many PPT charts store cached values in chart XML (c:numCache/c:strCache).
      - If not present, values might live in an embedded workbook; this keeps it lightweight.
    """
    path = _strip_pack_uri(chart_partname)
    if path not in zf.namelist():
        return {"source": "zip", "path": path, "series": [], "note": "chart xml not found in pptx zip"}

    xml_bytes = zf.read(path)
    root = ET.fromstring(xml_bytes)

    # Determine chart type by first matching chart element tag under c:plotArea
    chart_type = None
    plot_area = root.find(".//c:plotArea", NS)
    if plot_area is not None:
        for child in list(plot_area):
            tag = child.tag
            if tag.startswith("{%s}" % NS["c"]) and tag.endswith("Chart"):
                chart_type = tag.split("}", 1)[1]  # e.g. doughnutChart, barChart
                break

    series_out = []
    for ser in root.findall(".//c:ser", NS):
        # Series name
        sname = None
        tx = ser.find("./c:tx", NS)
        if tx is not None:
            # name could be in strRef/strCache or directly as v
            v = tx.find(".//c:v", NS)
            if v is not None and v.text is not None:
                sname = v.text

        # Categories (x-axis)
        cats = []
        cat = ser.find("./c:cat", NS)
        if cat is not None:
            for v in cat.findall(".//c:strCache//c:pt//c:v", NS):
                if v.text is not None:
                    cats.append(v.text)
            if not cats:
                for v in cat.findall(".//c:numCache//c:pt//c:v", NS):
                    if v.text is not None:
                        cats.append(v.text)

        # Values (y-axis)
        vals = []
        val = ser.find("./c:val", NS)
        if val is not None:
            for v in val.findall(".//c:numCache//c:pt//c:v", NS):
                if v.text is not None:
                    try:
                        vals.append(float(v.text))
                    except Exception:
                        vals.append(v.text)

        series_out.append({"name": sname, "categories": cats, "values": vals})

    return {"source": "zip", "path": path, "chart_type": chart_type, "series": series_out}


# ----------------------------
# Table extraction
# ----------------------------
def extract_table(shape) -> Dict[str, Any]:
    tbl = shape.table
    rows = tbl.rows
    cols = tbl.columns
    grid = []
    for r in range(len(rows)):
        row = []
        for c in range(len(cols)):
            cell = tbl.cell(r, c)
            row.append(cell.text if cell.text is not None else "")
        grid.append(row)
    return {
        "n_rows": len(rows),
        "n_cols": len(cols),
        "cells": grid,
    }


# ----------------------------
# Simple layout grouping heuristics (geometry-based)
# ----------------------------
def assign_columns(elements: List[Dict[str, Any]]) -> None:
    """
    Split into left/right columns using median of x-centers (normalized).
    Adds element['layout']['column'] = 'left'|'right'|'full'
    """
    xs = [e["bbox"]["center_norm"]["x"] for e in elements if "bbox" in e]
    if len(xs) < 6:
        for e in elements:
            e.setdefault("layout", {})["column"] = "full"
        return

    med = sorted(xs)[len(xs) // 2]
    for e in elements:
        cx = e["bbox"]["center_norm"]["x"]
        w = e["bbox"]["norm"]["w"]
        # if spans big portion, treat as full-width
        if w >= 0.80:
            col = "full"
        else:
            col = "left" if cx < med else "right"
        e.setdefault("layout", {})["column"] = col


def detect_title(elements: List[Dict[str, Any]]) -> Optional[int]:
    """
    Pick a likely slide title: top-most wide text with large font.
    Returns index of element in list or None.
    """
    candidates = []
    for i, e in enumerate(elements):
        if e["type"] not in ("textbox", "placeholder"):
            continue
        txt = (e.get("text") or "").strip()
        if not txt:
            continue
        y = e["bbox"]["norm"]["y"]
        w = e["bbox"]["norm"]["w"]
        font = e.get("max_font_pt") or 0.0
        score = 0.0
        score += max(0.0, 0.35 - y) * 3.0     # near top
        score += min(1.0, w / 0.9) * 2.0      # wide
        score += min(1.0, font / 32.0) * 2.0  # large font
        candidates.append((score, i))
    if not candidates:
        return None
    candidates.sort(reverse=True)
    return candidates[0][1]


# ----------------------------
# Main extraction
# ----------------------------
def extract_pptx(pptx_path: str) -> Dict[str, Any]:
    prs = Presentation(pptx_path)
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    with zipfile.ZipFile(pptx_path, "r") as zf:
        slides_out = []
        for sidx, slide in enumerate(prs.slides, start=1):
            elements = []

            # Preserve draw order (z): enumerate slide.shapes is typically back-to-front.
            for z, shape in enumerate(slide.shapes):
                etype = shape_type_name(shape)
                bbox = emu_box(shape, slide_w, slide_h)

                el: Dict[str, Any] = {
                    "id": f"s{sidx}_z{z}",
                    "type": etype,
                    "z": z,
                    "bbox": bbox,
                    "name": getattr(shape, "name", None),
                }

                # Text
                txt = safe_text(shape).strip()
                if txt:
                    el["text"] = txt
                    el["max_font_pt"] = max_font_pt(shape)

                # Table
                if etype == "table":
                    try:
                        el["table"] = extract_table(shape)
                    except Exception as e:
                        el["table_error"] = str(e)

                # Image
                if etype == "image":
                    try:
                        # Image bytes are in shape.image; we don't read pixels, only metadata.
                        img = shape.image
                        el["image"] = {
                            "filename": img.filename,
                            "content_type": img.content_type,
                            "size_bytes": len(img.blob) if img.blob is not None else None,
                        }
                    except Exception as e:
                        el["image_error"] = str(e)

                # Chart (try python-pptx first, then OOXML cache)
                if etype == "chart":
                    try:
                        chart = shape.chart
                        # chart.chart_type is an enum; str() gives readable-ish
                        el["chart"] = {"chart_type": str(chart.chart_type), "series": []}
                        for ser in chart.series:
                            sname = getattr(ser, "name", None)
                            vals = []
                            try:
                                # python-pptx exposes values as a sequence for many charts
                                vals = [v for v in ser.values]
                            except Exception:
                                vals = []
                            el["chart"]["series"].append({"name": sname, "values": vals})
                    except Exception as e:
                        el["chart_error"] = str(e)

                    # OOXML cached categories/values (more structured for pies/doughnuts)
                    try:
                        partname = str(shape.chart.part.partname)  # e.g. '/ppt/charts/chart1.xml'
                        el.setdefault("chart", {})
                        el["chart"]["ooxml_cache"] = parse_chart_xml_from_zip(zf, partname)
                    except Exception as e:
                        el.setdefault("chart", {})
                        el["chart"]["ooxml_cache_error"] = str(e)

                # Group: capture child shapes (basic)
                if etype == "group":
                    try:
                        kids = []
                        for kz, kshape in enumerate(shape.shapes):
                            kids.append({
                                "type": shape_type_name(kshape),
                                "bbox": emu_box(kshape, slide_w, slide_h),
                                "name": getattr(kshape, "name", None),
                                "text": safe_text(kshape).strip() or None
                            })
                        el["group_children"] = kids
                    except Exception as e:
                        el["group_error"] = str(e)

                elements.append(el)

            # Geometry-based enrichment
            assign_columns(elements)
            title_idx = detect_title(elements)
            if title_idx is not None:
                for i, e in enumerate(elements):
                    e.setdefault("layout", {})["is_title"] = (i == title_idx)

            slides_out.append({
                "slide_index": sidx,
                "slide_size_emu": {"w": slide_w, "h": slide_h},
                "elements": elements,
            })

    return {"file": os.path.basename(pptx_path), "slides": slides_out}


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx", help="Path to input .pptx")
    ap.add_argument("-o", "--out", default="pptx_extracted.json", help="Output JSON path")
    args = ap.parse_args()

    data = extract_pptx(args.pptx)
    with open(args.out, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

    print(f"Wrote: {args.out}")
    print("Note: If key slide content is embedded as images, you will not get text without OCR/vision.")


if __name__ == "__main__":
    main()