
from __future__ import annotations

from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Union

from pptx import Presentation


@dataclass
class ChartTable:
    slide_index: int
    shape_index: int
    shape_name: str
    chart_type: str
    table: List[List[Any]]          # 2D table suitable for CSV/Excel
    headers: List[str]              # convenience: first row of `table` (if present)


def extract_chart_data_tables(pptx_path: str) -> List[ChartTable]:
    """
    Extract chart *cached* data from a .pptx into simple 2D tables.

    What you get:
      - For category charts (column/line/bar/area, etc.): a table with columns:
            Category | <Series 1> | <Series 2> | ...
      - For scatter charts: a table with columns:
            Series | X | Y
      - For bubble charts (if present): Series | X | Y | Size

    Notes / limitations:
      - python-pptx does not expose the embedded Excel workbook as a first-class
        "read chart data table" API; this extracts the values visible via the
        chart object model (series values, categories, etc.). :contentReference[oaicite:0]{index=0}
      - If a chart is *linked* to external data, PPT may only store cached values;
        you'll get whatever is in the file.
    """
    prs = Presentation(pptx_path)
    out: List[ChartTable] = []

    for si, slide in enumerate(prs.slides):
        for shi, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_chart", False):
                continue

            chart = shape.chart
            chart_type = str(getattr(chart, "chart_type", "UNKNOWN"))

            # Try category plot first (common case): plot.categories + series.values :contentReference[oaicite:1]{index=1}
            table: List[List[Any]] = []
            headers: List[str] = []

            extracted = False

            # Charts can have multiple plots; we merge series across plots by using chart.series :contentReference[oaicite:2]{index=2}
            series_list = list(getattr(chart, "series", []))

            # Attempt to pull categories from the first plot when available.
            categories: Optional[List[Any]] = None
            try:
                plots = list(chart.plots)  # may raise on some chart types
                if plots:
                    plot0 = plots[0]
                    if hasattr(plot0, "categories"):
                        categories = [c.label for c in plot0.categories]
            except Exception:
                categories = None

            # CATEGORY-BASED TABLE
            if categories is not None and series_list:
                headers = ["Category"] + [getattr(s, "name", f"Series {i+1}") for i, s in enumerate(series_list)]
                table.append(headers)

                # Build rows: each category index corresponds to one value per series
                n = len(categories)
                for idx in range(n):
                    row = [categories[idx]]
                    for s in series_list:
                        # s.values is a sequence of numbers :contentReference[oaicite:3]{index=3}
                        vals = list(getattr(s, "values", []))
                        row.append(vals[idx] if idx < len(vals) else None)
                    table.append(row)

                extracted = True

            # SCATTER / XY TABLE
            if not extracted and series_list:
                rows = [["Series", "X", "Y"]]
                ok = False
                for s in series_list:
                    # XYSeries typically exposes x_values and y_values; not present on category charts.
                    x_vals = getattr(s, "x_values", None)
                    y_vals = getattr(s, "y_values", None)
                    if x_vals is None or y_vals is None:
                        continue
                    x_vals = list(x_vals)
                    y_vals = list(y_vals)
                    name = getattr(s, "name", "Series")
                    m = min(len(x_vals), len(y_vals))
                    for i in range(m):
                        rows.append([name, x_vals[i], y_vals[i]])
                    ok = True

                if ok:
                    table = rows
                    headers = rows[0]
                    extracted = True

            # BUBBLE (best-effort)
            if not extracted and series_list:
                rows = [["Series", "X", "Y", "Size"]]
                ok = False
                for s in series_list:
                    x_vals = getattr(s, "x_values", None)
                    y_vals = getattr(s, "y_values", None)
                    sizes = getattr(s, "bubble_sizes", None)  # may not exist
                    if x_vals is None or y_vals is None or sizes is None:
                        continue
                    x_vals = list(x_vals)
                    y_vals = list(y_vals)
                    sizes = list(sizes)
                    name = getattr(s, "name", "Series")
                    m = min(len(x_vals), len(y_vals), len(sizes))
                    for i in range(m):
                        rows.append([name, x_vals[i], y_vals[i], sizes[i]])
                    ok = True

                if ok:
                    table = rows
                    headers = rows[0]
                    extracted = True

            # If we couldn't extract in a structured way, still record presence.
            if not extracted:
                table = [["Unsupported chart type (extracted no series/categories via python-pptx).", chart_type]]
                headers = table[0]

            out.append(
                ChartTable(
                    slide_index=si,
                    shape_index=shi,
                    shape_name=getattr(shape, "name", f"shape_{shi}"),
                    chart_type=chart_type,
                    table=table,
                    headers=headers,
                )
            )

    return out
